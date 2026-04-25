"""Generate a detailed research presentation from project results."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme Colors ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x25, 0x25, 0x3D)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_GREEN = RGBColor(0x7B, 0xC6, 0x7E)
ACCENT_ORANGE = RGBColor(0xE8, 0xA8, 0x38)
ACCENT_RED = RGBColor(0xD9, 0x4A, 0x4A)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
TABLE_HEADER_BG = RGBColor(0x4A, 0x90, 0xD9)
TABLE_ROW_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, left=0.5, top=0.3, width=9, height=0.8,
                  font_size=28, color=WHITE, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def add_body_text(slide, text, left=0.5, top=1.3, width=9, height=5.5,
                  font_size=14, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf


def add_bullet_slide(slide, title, bullets, sub_bullets=None):
    set_slide_bg(slide)
    add_title_box(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(6)
        p.level = 0

        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                p2 = tf.add_paragraph()
                p2.text = sub
                p2.font.size = Pt(13)
                p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
                p2.space_after = Pt(3)
                p2.level = 1


def add_table(slide, data, left=0.3, top=1.4, width=9.4, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.35 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.color.rgb = DARK_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_LIGHT if r % 2 == 0 else TABLE_ROW_WHITE

    return table


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ════════════════════════════════════════════
    # SLIDE 1: Title
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Evaluating Self-Reflection and Error Correction\nin Agent-Based Defensive Security Threat Detection",
                  top=1.5, font_size=28)
    add_body_text(slide,
                  "Research Question:\nDoes incorporating self-reflection and error correction significantly\nimprove the accuracy and reliability of agent-based defensive\nsecurity threat detection systems?",
                  top=3.2, font_size=16, color=ACCENT_BLUE)
    add_body_text(slide,
                  "Evaluated across 4 threat domains | 3 reflection levels | 3 Gemini models",
                  top=5.2, font_size=14, color=LIGHT_GRAY)

    # ════════════════════════════════════════════
    # SLIDE 2: Introduction
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Introduction", [
        "LLMs show promise in cybersecurity: phishing detection, malware classification, log analysis",
        "Problem: Single-pass LLM analysis is prone to errors",
        "False Positives: Waste analyst time, cause alert fatigue",
        "False Negatives: Miss real threats, lead to breaches",
        "Solution: Self-reflection — agent reviews and revises its own reasoning",
        "Mirrors human workflow: junior analyst → senior review → revised assessment",
        "This project: 3-level reflection architecture across 4 security domains",
    ])

    # ════════════════════════════════════════════
    # SLIDE 3: What is Self-Reflection
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "What is Self-Reflection?")
    add_body_text(slide, (
        "Self-reflection in AI agents = reviewing, critiquing, and revising own output\n\n"
        "Three roles played by a single LLM through different prompts:\n\n"
        "  1. Detection Agent — performs initial threat analysis\n"
        "  2. Critic Agent — reviews for errors, missed indicators, logical gaps\n"
        "  3. Revision Step — reconsiders verdict using critic's feedback\n\n"
        "Three Levels of Reflection:\n\n"
        "  Level 0 (Baseline):   Sample ──> Agent ──> Verdict\n"
        "  Level 1 (Single):     Sample ──> Agent ──> Critic ──> Revise ──> Verdict\n"
        "  Level 2 (Iterative):  Sample ──> Agent <──> Critic (loop max 3x) ──> Verdict"
    ), font_size=14)

    # ════════════════════════════════════════════
    # SLIDE 4: Threat Domains Overview
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "4 Threat Domains")
    add_table(slide, [
        ["Domain", "Dataset", "Samples", "Input Type", "Class Balance"],
        ["Phishing Email", "CEAS_08", "39,154", "Email subject, sender, body", "56% phishing / 44% benign"],
        ["Network Intrusion", "NSL-KDD", "125,972", "41 numeric features", "Mixed attacks + normal"],
        ["Malware (PE)", "ClaMP", "5,210", "70 PE header features", "52% malware / 48% benign"],
        ["Insider Threat Logs", "CERT r4.2", "242,228", "Window of 7 events", "5.3% threat / 94.7% benign"],
    ])

    # ════════════════════════════════════════════
    # SLIDE 5: Phishing Domain
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Domain 1: Phishing Email Detection", [
        "Dataset: CEAS_08 — 39,154 emails (21,842 phishing + 17,312 benign)",
        "Input: Email subject, sender address, and body text",
        "Challenge: Distinguishing phishing from legitimate marketing, newsletters, notifications",
        "Key Indicators: Spoofed domains, urgency language, suspicious URLs, social engineering",
        "LLM Strength: Excels here — natural language understanding is core LLM capability",
        "Results: 93-100% accuracy across models at Level 0",
    ])

    # ════════════════════════════════════════════
    # SLIDE 6: Network Domain
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Domain 2: Network Intrusion Detection", [
        "Dataset: NSL-KDD — 125,972 connection records with 41 features each",
        "Input: Structured features (duration, protocol, bytes, error rates, etc.)",
        "Challenge: Classify traffic as normal or attack (DoS, Probe, R2L, U2R)",
        "Key Indicators: Traffic anomalies, SYN floods, port scans, failed logins",
        "LLM Strength: Good — recognizable attack patterns in feature combinations",
        "Results: 90-100% accuracy across models at Level 0",
    ])

    # ════════════════════════════════════════════
    # SLIDE 7: Malware Domain
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Domain 3: Malware Analysis (PE Features)", [
        "Dataset: ClaMP — 5,210 PE files with 70 features each",
        "Input: Raw PE header values (entropy, sections, linker version, flags...)",
        "Challenge: Classify Windows executables from numerical features only",
        "Key Indicators: Entropy patterns, suspicious sections, packer detection",
        "LLM Weakness: Cannot reason about 70 raw numbers — scores 50% (random chance)",
        "All models default to calling everything 'malicious' (FPR = 1.0)",
    ])

    # ════════════════════════════════════════════
    # SLIDE 8: Log Domain
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Domain 4: Insider Threat Log Detection", [
        "Dataset: CERT r4.2 — 242,228 events (94.7% benign, 5.3% malicious)",
        "Input: Sliding window of 7 consecutive employee events",
        "Challenge: Malicious and benign windows look nearly identical",
        "USB Connect → File Access → File Access appears in BOTH classes",
        "Missing Context: Employee role, resignation status, job scope, time of day",
        "Labels based on WHO the user is, not WHAT the events look like",
        "Results: 37-50% accuracy — near random chance across all models",
    ])

    # ════════════════════════════════════════════
    # SLIDE 9: Tech Stack
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Tech Stack")
    add_table(slide, [
        ["Component", "Technology", "Purpose"],
        ["LLM Provider", "Google Gemini API", "Cloud-hosted language models"],
        ["Agent Framework", "LangChain", "Prompt management, chain composition"],
        ["Prompt Templates", "ChatPromptTemplate", "Structured prompts per agent role"],
        ["Chain Composition", "LCEL (pipe operator)", "prompt | llm pipeline"],
        ["Output Parsing", "Manual JSON + validation", "Robust LLM response handling"],
        ["Evaluation", "scikit-learn, scipy", "Metrics, McNemar's test"],
        ["Data Processing", "pandas", "Loading, preprocessing, windowing"],
        ["Visualization", "matplotlib, seaborn, plotly", "Charts and figures"],
        ["Demo", "Streamlit", "Interactive web demo"],
        ["Config", "python-dotenv", "API key management"],
    ], top=1.2)

    # ════════════════════════════════════════════
    # SLIDE 10: Models Tested
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Models Tested")
    add_table(slide, [
        ["Model", "Type", "Input Cost", "Output Cost", "Key Trait"],
        ["gemini-2.5-flash-lite", "Lightweight", "$0.075/1M", "$0.30/1M", "Fast, cheap, no thinking tokens"],
        ["gemini-2.5-flash", "Hybrid reasoning", "$0.15/1M", "$0.60 + $3.50 thinking/1M", "Thinking tokens for deep reasoning"],
        ["gemini-3-flash-preview", "Frontier", "$0.50/1M", "$3.00/1M", "Most capable reasoning"],
        ["Ollama (Llama/Mistral)", "Open-source local", "Free", "Free", "Failed: can't produce valid JSON"],
    ])

    # ════════════════════════════════════════════
    # SLIDE 11: System Architecture
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "System Architecture")
    add_body_text(slide, (
        "┌─────────────┐     ┌──────────────────┐     ┌─────────────────┐\n"
        "│   Datasets   │────>│    DataLoader     │────>│ ReflectiveAgent │\n"
        "│  CEAS_08     │     │  load, balance,   │     │  (Orchestrator) │\n"
        "│  NSL-KDD     │     │  format, window   │     │                 │\n"
        "│  ClaMP       │     └──────────────────┘     └────────┬────────┘\n"
        "│  CERT r4.2   │                                       │\n"
        "└─────────────┘                           ┌────────────┼────────────┐\n"
        "                                          v            v            v\n"
        "┌──────────────┐              ┌───────────────┐ ┌────────────┐ ┌──────────┐\n"
        "│ Domain       │──────────────>│  Detection    │ │   Critic   │ │ Revision │\n"
        "│ Prompts      │              │  Agent (L0)   │ │   Agent    │ │   Step   │\n"
        "│ phishing.py  │              └───────┬───────┘ └─────┬──────┘ └────┬─────┘\n"
        "│ network.py   │                      │               │             │\n"
        "│ malware.py   │                      └───────────────┴─────────────┘\n"
        "│ log.py       │                                      │\n"
        "└──────────────┘                              ┌───────v───────┐\n"
        "                                              │  Gemini API   │\n"
        "                                              │  (LangChain)  │\n"
        "                                              └───────────────┘"
    ), font_size=10)

    # ════════════════════════════════════════════
    # SLIDE 12: Level 0 Explained
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Level 0: Baseline (No Reflection)", [
        "Flow: Sample ──> Detection Agent ──> Final Verdict",
        "LLM Calls: 1 per sample",
        "How it works:",
        "  1. Sample text sent to LLM with domain-specific system prompt",
        "  2. LLM returns JSON: verdict, confidence, reasoning, indicators, threat_type",
        "  3. That's the final answer — no review, no revision",
        "LangChain: ChatPromptTemplate.from_messages() | ChatGoogleGenerativeAI",
        "This is the CONTROL GROUP — all improvements measured against Level 0",
    ])

    # ════════════════════════════════════════════
    # SLIDE 13: Level 1 Explained
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Level 1: Single Reflection", [
        "Flow: Sample ──> Detection ──> Critic ──> Revision ──> Final Verdict",
        "LLM Calls: 3 per sample (always)",
        "Step 1: Detection Agent produces initial verdict (same as Level 0)",
        "Step 2: Critic Agent receives sample + full analysis, reviews for errors",
        "Step 3: Revision with context-aware instructions:",
    ], sub_bullets={
        4: [
            "If critic AGREED → Keep verdict, refine reasoning and confidence",
            "If critic same verdict, different reasoning → Improve reasoning only",
            "If critic DISAGREED → Compare evidence both sides, keep if yours is stronger",
        ]
    })

    # ════════════════════════════════════════════
    # SLIDE 14: Level 2 Explained
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Level 2: Iterative Reflection (Max 3 Rounds)", [
        "Flow: Sample ──> Detection ──> [Critic ──> Revise] x N ──> Final Verdict",
        "LLM Calls: 2-7 per sample (depends on consensus)",
        "Loop repeats until CONSENSUS or max 3 rounds:",
        "Consensus requires ALL THREE conditions:",
    ], sub_bullets={
        3: [
            "Critic agrees with current verdict (agree: true)",
            "Critic's verdict MATCHES Detection Agent's verdict",
            "Both confidences >= 0.7 (consensus threshold)",
        ]
    })
    add_body_text(slide,
                  "On consensus: final confidence = average(agent, critic)\n"
                  "No consensus after 3 rounds: uses last revised verdict",
                  top=5.5, font_size=13, color=ACCENT_ORANGE)

    # ════════════════════════════════════════════
    # SLIDE 15: LangChain Implementation
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "LangChain Implementation Details", [
        "ChatGoogleGenerativeAI — LangChain's Gemini integration, temperature=0",
        "ChatPromptTemplate.from_messages() — structured prompts for each agent role",
        "LCEL Chain: chain = prompt | llm  →  response = chain.invoke({...})",
        "Manual JSON Parsing — handles ```json blocks, validates fields, clamps confidence",
        "Pydantic Models — ThreatVerdict and CriticReview define expected output schemas",
        "TokenCountingCallback — LangChain callback for tracking token usage per call",
        "Timed Calls — _timed_call() wrapper records latency per LLM call (detection, critic, revision)",
        "Domain Context Caching — _get_domain_context() cached at agent init, not per sample",
    ])

    # ════════════════════════════════════════════
    # SLIDE 16: Evaluation Metrics
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Evaluation Metrics")
    add_table(slide, [
        ["Metric", "Description", "What It Tells Us"],
        ["Accuracy", "Correct / Total", "Overall performance"],
        ["Precision", "TP / (TP + FP)", "When it says malicious, is it right?"],
        ["Recall", "TP / (TP + FN)", "Does it catch all threats?"],
        ["F1 Score", "Harmonic mean of P & R", "Balanced single metric"],
        ["FPR", "FP / (FP + TN)", "Benign flagged as threats"],
        ["FNR", "FN / (FN + TP)", "Threats missed as benign"],
        ["McNemar's Test", "Chi-squared on discordant pairs", "Is L0 vs L2 difference significant?"],
    ])

    # ════════════════════════════════════════════
    # SLIDE 17: Results Flash-Lite
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Results: Gemini 2.5 Flash-Lite (30 samples)", font_size=24)
    add_table(slide, [
        ["Domain", "Level", "Accuracy", "Precision", "Recall", "F1", "FPR", "FNR"],
        ["Phishing", "L0", "0.9333", "0.8824", "1.0000", "0.9375", "0.1333", "0.0000"],
        ["", "L1", "0.9333", "1.0000", "0.8667", "0.9286", "0.0000", "0.1333"],
        ["", "L2", "0.8000", "1.0000", "0.6000", "0.7500", "0.0000", "0.4000"],
        ["Network", "L0", "0.9000", "0.8750", "0.9333", "0.9032", "0.1333", "0.0667"],
        ["", "L1", "0.8000", "0.9091", "0.6667", "0.7692", "0.0667", "0.3333"],
        ["", "L2", "0.7333", "0.8889", "0.5333", "0.6667", "0.0667", "0.4667"],
        ["Malware", "L0", "0.5000", "0.5000", "1.0000", "0.6667", "1.0000", "0.0000"],
        ["", "L1", "0.3667", "0.3000", "0.2000", "0.2400", "0.4667", "0.8000"],
        ["", "L2", "0.5000", "0.5000", "0.2667", "0.3478", "0.2667", "0.7333"],
        ["Logs", "L0", "0.3667", "0.3571", "0.3333", "0.3448", "0.6000", "0.6667"],
        ["", "L1", "0.5000", "0.0000", "0.0000", "0.0000", "0.0000", "1.0000"],
        ["", "L2", "0.5000", "0.0000", "0.0000", "0.0000", "0.0000", "1.0000"],
    ], top=1.2)

    # ════════════════════════════════════════════
    # SLIDE 18: Flash-Lite Reflection Impact
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Flash-Lite: Reflection Impact (L0 → L2)", font_size=24)
    add_table(slide, [
        ["Domain", "L0 Acc", "L2 Acc", "Change", "L0 F1", "L2 F1", "Change"],
        ["Phishing", "0.9333", "0.8000", "↓ 0.1333", "0.9375", "0.7500", "↓ 0.1875"],
        ["Network", "0.9000", "0.7333", "↓ 0.1667", "0.9032", "0.6667", "↓ 0.2366"],
        ["Malware", "0.5000", "0.5000", "= 0.0000", "0.6667", "0.3478", "↓ 0.3188"],
        ["Logs", "0.3667", "0.5000", "↑ 0.1333", "0.3448", "0.0000", "↓ 0.3448"],
    ])
    add_body_text(slide,
                  "Key Finding: With flash-lite, self-reflection consistently DEGRADES performance.\n"
                  "The model is too weak for multi-step reasoning — the critic second-guesses correct\n"
                  "answers, and the revision agent caves to the critic's suggestions.",
                  top=4.5, font_size=14, color=ACCENT_RED)

    # ════════════════════════════════════════════
    # SLIDE 19: Results Flash
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Results: Gemini 2.5 Flash (6 samples)", font_size=24)
    add_table(slide, [
        ["Domain", "Level", "Accuracy", "Precision", "Recall", "F1", "FPR", "FNR"],
        ["Phishing", "L0-L2", "1.0000", "1.0000", "1.0000", "1.0000", "0.0000", "0.0000"],
        ["Network", "L0-L2", "1.0000", "1.0000", "1.0000", "1.0000", "0.0000", "0.0000"],
        ["Malware", "L0-L2", "0.5000", "0.5000", "1.0000", "0.6667", "1.0000", "0.0000"],
        ["Logs", "L0-L2", "0.5000", "0.5000", "1.0000", "0.6667", "1.0000", "0.0000"],
    ])
    add_body_text(slide,
                  "Key Finding: Ceiling effect — perfect on text domains, reflection has nothing\n"
                  "to improve. Still fails on malware/logs (50% = random chance).",
                  top=4.0, font_size=14, color=ACCENT_ORANGE)

    # ════════════════════════════════════════════
    # SLIDE 20: Results Flash-3
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Results: Gemini 3 Flash Preview (10 samples)", font_size=24)
    add_table(slide, [
        ["Domain", "Level", "Accuracy", "Precision", "Recall", "F1", "FPR", "FNR"],
        ["Phishing", "L0-L2", "1.0000", "1.0000", "1.0000", "1.0000", "0.0000", "0.0000"],
        ["Network", "L0-L2", "1.0000", "1.0000", "1.0000", "1.0000", "0.0000", "0.0000"],
        ["Malware", "L0-L2", "0.5000", "0.5000", "1.0000", "0.6667", "1.0000", "0.0000"],
        ["Logs", "L0", "0.4000", "0.4444", "0.8000", "0.5714", "1.0000", "0.2000"],
        ["", "L1", "0.2000", "0.2000", "0.2000", "0.2000", "0.8000", "0.8000"],
        ["", "L2", "0.3000", "0.3333", "0.4000", "0.3636", "0.8000", "0.6000"],
    ])
    add_body_text(slide,
                  "Key Finding: Most capable model confirms pattern — perfect on text,\n"
                  "fails on PE features. On logs, reflection HURTS (40% → 30%).",
                  top=5.0, font_size=14, color=ACCENT_RED)

    # ════════════════════════════════════════════
    # SLIDE 21: Statistical Significance
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Statistical Significance (McNemar's Test)")
    add_body_text(slide,
                  "McNemar's Test compares Level 0 vs Level 2 predictions on the same samples.\n"
                  "Tests whether the classifiers have significantly different error rates.\n",
                  top=1.3, font_size=14)
    add_table(slide, [
        ["Model", "Phishing", "Network", "Malware", "Logs"],
        ["Flash-Lite (30)", "p=0.2891 NO", "p=0.1250 NO", "p=1.0000 NO", "p=0.4240 NO"],
        ["Flash (6)", "p=1.0000 NO", "p=1.0000 NO", "p=1.0000 NO", "p=1.0000 NO"],
        ["Flash-3 (10)", "p=1.0000 NO", "p=1.0000 NO", "p=1.0000 NO", "p=1.0000 NO"],
    ], top=3.0)
    add_body_text(slide,
                  "No statistically significant differences (all p > 0.05).\n"
                  "Partly due to small sample sizes, partly because differences were not large enough.",
                  top=5.2, font_size=13, color=LIGHT_GRAY)

    # ════════════════════════════════════════════
    # SLIDE 22: Model Comparison
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Model Comparison: Capability vs Reflection", font_size=24)
    add_table(slide, [
        ["Model", "Phishing L0", "Network L0", "Reflection Helps?", "Why"],
        ["flash-lite", "93.3%", "90.0%", "NO — hurts", "Too weak to be good critic"],
        ["2.5-flash", "100%", "100%", "Neutral — ceiling", "Already perfect"],
        ["3-flash-preview", "100%", "100%", "Neutral — ceiling", "Already perfect"],
        ["Ollama (local)", "N/A", "N/A", "N/A — broken", "Can't produce JSON"],
    ])
    add_body_text(slide,
                  "The Self-Reflection Paradox:\n"
                  "  - Weak models: Critic gives bad feedback → revision caves → accuracy drops\n"
                  "  - Strong models: Already perfect → nothing to improve → reflection is redundant\n"
                  "  - Hard domains: No model can do it → reflection adds noise to random guessing",
                  top=4.5, font_size=13, color=ACCENT_ORANGE)

    # ════════════════════════════════════════════
    # SLIDE 23: Ollama Failure
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Why Open-Source / Ollama Models Failed")
    add_table(slide, [
        ["Model", "JSON Output", "Critic Quality", "Verdict"],
        ["Llama 3.1 (8B)", "Frequently malformed", "N/A — parse failed", "FAILED"],
        ["Mistral (7B)", "Inconsistent formatting", "N/A — parse failed", "FAILED"],
        ["Gemma 2 (9B)", "Partial, often truncated", "N/A — parse failed", "FAILED"],
    ], top=1.4)
    add_bullet_slide(slide, "", [])
    add_body_text(slide, (
        "The self-reflection pipeline requires EVERY LLM call to return valid JSON.\n"
        "One parse failure at any step (detection, critic, revision) breaks the chain.\n\n"
        "What it would take to use open-source models:\n"
        "  - 70B+ parameter models (Llama 3.1 70B) — need 40GB+ VRAM\n"
        "  - Fine-tuning smaller models on structured JSON output\n"
        "  - Capable GPU hardware (~₹1-2 lakh investment)\n"
        "  - 'Free' local models actually require significant compute cost"
    ), top=3.5, font_size=13)

    # ════════════════════════════════════════════
    # SLIDE 24: Malware Problem
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "The Malware Domain Problem", [
        "Every model scores exactly 50% — random chance",
        "All models default to calling everything 'malicious' (FPR = 1.0)",
        "Why: ClaMP has 70 raw numeric PE features (e_cblp: 144, SizeOfCode: 22528...)",
        "LLMs are LANGUAGE models, not NUMERICAL classifiers",
        "They cannot learn statistical distributions from a single prompt",
        "They cannot perform feature correlation across 70 dimensions",
        "Traditional ML (Random Forest, XGBoost) achieves 95%+ easily on same data",
        "Solution: Hybrid approach — ML extracts features, LLM interprets in natural language",
    ])

    # ════════════════════════════════════════════
    # SLIDE 25: Log Problem
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "The Insider Threat Log Problem")
    add_body_text(slide, (
        "Malicious Window:                        Benign Window:\n"
        "─────────────────                        ─────────────────\n"
        "Event 1: device Connect | 10:51          Event 1: device Connect | 11:45\n"
        "Event 2: file PU444C7E.doc | 10:53       Event 2: file report.doc | 11:47\n"
        "Event 3: file OEDCGST5.pdf | 10:57       Event 3: file data.xlsx | 11:52\n"
        "Event 4: file GY2Y0OPE.doc | 11:00       Event 4: device Disconnect | 11:55\n\n"
        "They look IDENTICAL. Both show: USB connect → file access → file access\n\n"
        "The difference is CONTEXT that doesn't exist in the log events:\n"
        "  - Has the employee submitted a resignation?\n"
        "  - Are these files outside their job scope?\n"
        "  - Is this happening at 2 AM or 2 PM?\n"
        "  - What is the employee's role and department?\n\n"
        "CERT labels are based on WHO + WHEN, not WHAT the events look like."
    ), font_size=13)

    # ════════════════════════════════════════════
    # SLIDE 26: Tradeoff Triangle
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "The Tradeoff Triangle")
    add_body_text(slide, (
        "                    ACCURACY\n"
        "                       ^\n"
        "                      /|\\\n"
        "                     / | \\\n"
        "                    /  |  \\\n"
        "                   / gemini \\\n"
        "                  / 3-flash  \\\n"
        "                 / (100% but  \\\n"
        "                /  ceiling)    \\\n"
        "               /      |        \\\n"
        "              / \"Goldilocks\"    \\\n"
        "             /  Zone — where     \\\n"
        "            /  reflection COULD   \\\n"
        "           /     help, but we      \\\n"
        "          /   couldn't find a       \\\n"
        "         /  model that sits here     \\\n"
        "        /         |                  \\\n"
        "       / flash-lite    Ollama         \\\n"
        "      / (93% but     (can't even      \\\n"
        "     / reflection    produce JSON)     \\\n"
        "    /   hurts)                          \\\n"
        "   +────────────+───────────────────────>\n"
        "  FREE        $0.10                $2.60   COST"
    ), font_size=11)

    # ════════════════════════════════════════════
    # SLIDE 27: Cost Analysis
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Cost Analysis")
    add_table(slide, [
        ["Model", "Input/1M", "Output/1M", "30 samples est. (USD)", "INR"],
        ["gemini-2.5-flash-lite", "$0.075", "$0.30", "~$0.10", "~₹10"],
        ["gemini-2.5-flash", "$0.15", "$0.60 + $3.50 think", "~$2.00", "~₹170"],
        ["gemini-3-flash-preview", "$0.50", "$3.00", "~$2.60", "~₹220"],
    ], top=1.4)
    add_body_text(slide, "\nCost per Reflection Level (relative):", top=3.5, font_size=14, color=WHITE)
    add_table(slide, [
        ["Level", "LLM Calls/Sample", "Relative Cost", "Accuracy Impact"],
        ["L0 (Baseline)", "1", "1x", "Best or tied"],
        ["L1 (Single)", "3", "3x", "Same or worse"],
        ["L2 (Iterative)", "2-7", "2-7x", "Same or worse"],
    ], top=4.2)

    # ════════════════════════════════════════════
    # SLIDE 28: Key Findings
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, "Key Findings", [
        "1. Self-reflection requires a MINIMUM model capability threshold",
        "   Flash-lite lacks reasoning depth — critic gives bad advice, revision follows blindly",
        "2. Domain complexity determines reflection potential",
        "   Text domains: LLMs excel → ceiling effect. Numeric domains: LLMs fail → reflection can't help",
        "3. Stronger models hit ceiling effects",
        "   gemini-2.5-flash and 3-flash achieve 100% on phishing/network with just 6-10 samples",
        "4. Reflection adds significant cost with marginal or negative returns",
        "   L2 costs 2-7x more than L0 but never outperformed it in our experiments",
        "5. No statistical significance achieved",
        "   All McNemar's test p-values > 0.05 across all domain-model combinations",
    ])

    # ════════════════════════════════════════════
    # SLIDE 29: Project Structure
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Project Structure")
    add_body_text(slide, (
        "project/\n"
        "├── src/\n"
        "│   ├── agents/\n"
        "│   │   ├── base_agent.py          # Level 0: ThreatDetectionAgent\n"
        "│   │   ├── critic_agent.py        # Critic: CriticAgent\n"
        "│   │   └── reflective_agent.py    # Orchestrator: Levels 0/1/2\n"
        "│   ├── threats/\n"
        "│   │   ├── phishing.py            # Phishing prompts + formatting\n"
        "│   │   ├── network_intrusion.py   # Network prompts + formatting\n"
        "│   │   ├── malware.py             # Malware prompts + formatting\n"
        "│   │   └── log_analysis.py        # Log prompts + formatting\n"
        "│   ├── data/\n"
        "│   │   └── loader.py              # DataLoader: CSV, windowing, balancing\n"
        "│   ├── evaluation/\n"
        "│   │   ├── metrics.py             # Accuracy, F1, McNemar's test\n"
        "│   │   └── cost_tracker.py        # Token & latency tracking\n"
        "│   └── config.py                  # Model, thresholds, paths\n"
        "├── experiments/\n"
        "│   └── run_experiment.py          # CLI experiment runner\n"
        "├── data/                          # Datasets (gitignored)\n"
        "├── notebooks/                     # Analysis notebooks\n"
        "└── app/streamlit_app.py           # Interactive demo"
    ), font_size=11)

    # ════════════════════════════════════════════
    # SLIDE 30: Conclusion
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_box(slide, "Conclusion", font_size=32)
    add_body_text(slide, (
        "Self-reflection is NOT a universal improvement.\n\n"
        "It requires careful matching of:\n\n"
        "    Model Capability  ×  Domain Complexity  ×  Dataset Difficulty\n\n\n"
        "The ideal operating point — where the model is good enough to benefit\n"
        "from reflection but not so good that it's already perfect —\n"
        "remains elusive with current models and datasets.\n\n\n"
        "This is itself a valuable research finding:\n"
        "Self-reflection adds cost and complexity that is only justified when\n"
        "the base model operates in a narrow 'Goldilocks zone' of capability."
    ), font_size=16, color=LIGHT_GRAY)

    # ════════════════════════════════════════════
    # SAVE
    # ════════════════════════════════════════════
    output_path = "LLM_Self_Reflection_Security_Threat_Detection.pptx"
    prs.save(output_path)
    print(f"Presentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
